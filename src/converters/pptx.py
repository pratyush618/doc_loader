from typing import Optional, BinaryIO
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..core.exceptions import ConversionException
from .base import BaseConverter, ConverterResult


class PPTXConverter(BaseConverter):
    """Converter for Microsoft PowerPoint PPTX files"""
    
    def accepts(self, file: BinaryIO, mimetype: Optional[str] = None,
                extension: Optional[str] = None) -> bool:
        """Check if this is a PPTX file"""
        if extension and extension.lower() in {'.pptx', '.pptm', '.potx', '.potm'}:
            return True
        
        if mimetype and mimetype in {
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint.presentation.macroEnabled.12',
            'application/vnd.openxmlformats-officedocument.presentationml.template',
            'application/vnd.ms-powerpoint.template.macroEnabled.12'
        }:
            return True
        
        # Try to open with python-pptx
        try:
            Presentation(file)
            return True
        except (PackageNotFoundError, Exception):
            return False
        finally:
            self._reset_file_position(file)
    
    async def convert(self, file: BinaryIO, **kwargs) -> ConverterResult:
        """Convert PPTX to markdown"""
        try:
            # Load presentation
            prs = Presentation(file)
            
            # Convert slides to markdown
            markdown_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"# Slide {slide_num}\n")
                
                # Extract text from all shapes
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Try to detect if this is a title based on position and size
                        if hasattr(shape, 'top') and hasattr(shape, 'height'):
                            # If shape is in upper portion, treat as heading
                            if shape.top < prs.slide_height * 0.3:
                                slide_content.append(f"## {text}")
                            else:
                                # Check if it looks like a bullet point
                                lines = text.split('\n')
                                formatted_lines = []
                                for line in lines:
                                    line = line.strip()
                                    if line:
                                        # Check for bullet points or numbered lists
                                        if line.startswith(('•', '-', '*')) or any(line.startswith(f"{i}.") for i in range(1, 10)):
                                            formatted_lines.append(f"- {line.lstrip('•-* ').lstrip('0123456789. ')}")
                                        else:
                                            formatted_lines.append(line)
                                
                                if formatted_lines:
                                    slide_content.extend(formatted_lines)
                        else:
                            slide_content.append(text)
                    
                    # Handle tables
                    elif hasattr(shape, "table"):
                        table_md = self._convert_table_to_markdown(shape.table)
                        if table_md:
                            slide_content.append(table_md)
                
                if slide_content:
                    markdown_parts.append("\n".join(slide_content))
                else:
                    markdown_parts.append("*No text content*")
                
                markdown_parts.append("\n")
            
            content = "\n".join(markdown_parts)
            
            # Extract title from first slide
            title = None
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        title = shape.text.strip()[:100]  # First 100 chars
                        break
            
            if not title:
                title = f"PowerPoint Presentation ({len(prs.slides)} slides)"
            
            return ConverterResult(
                content=content,
                title=title
            )
            
        except Exception as e:
            raise ConversionException(f"Failed to convert PPTX: {str(e)}")
    
    def _convert_table_to_markdown(self, table) -> str:
        """Convert PowerPoint table to markdown"""
        if not table.rows:
            return ""
        
        markdown_lines = []
        
        # Process rows
        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
                row_cells.append(cell_text)
            
            markdown_lines.append("| " + " | ".join(row_cells) + " |")
            
            # Add separator after first row (header)
            if row_idx == 0:
                separators = ["-" * max(len(cell), 3) for cell in row_cells]
                markdown_lines.append("| " + " | ".join(separators) + " |")
        
        return "\n".join(markdown_lines) if markdown_lines else ""